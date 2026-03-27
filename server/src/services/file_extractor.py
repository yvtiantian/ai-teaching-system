"""文件文本提取服务 — 支持 PDF / DOCX / PPTX / TXT"""

from __future__ import annotations

import io

from loguru import logger

# 提取后文本最大字符数（约 12k tokens）
MAX_TEXT_LENGTH = 50_000

# 支持的 MIME → 提取函数映射（懒加载）
_SUPPORTED_MIMES: dict[str, str] = {
    "application/pdf": "_extract_pdf",
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document": "_extract_docx",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation": "_extract_pptx",
    "text/plain": "_extract_txt",
}


def supported_mime_types() -> list[str]:
    """返回当前支持的 MIME 类型列表。"""
    return list(_SUPPORTED_MIMES.keys())


def extract_text(file_bytes: bytes, mime_type: str, *, max_length: int = MAX_TEXT_LENGTH) -> str:
    """统一入口：根据 MIME 类型提取文件文本。

    Args:
        file_bytes: 文件二进制内容
        mime_type: MIME 类型
        max_length: 最大返回文本长度

    Returns:
        提取的纯文本

    Raises:
        ValueError: 不支持的 MIME 类型
        RuntimeError: 提取过程出错
    """
    handler_name = _SUPPORTED_MIMES.get(mime_type)
    if handler_name is None:
        raise ValueError(f"不支持的文件类型: {mime_type}")

    handler = globals()[handler_name]
    try:
        text: str = handler(file_bytes)
    except Exception as exc:
        logger.error("文件提取失败 ({}): {}", mime_type, exc)
        raise RuntimeError(f"文件内容提取失败: {exc}") from exc

    # 截断
    if len(text) > max_length:
        text = text[:max_length] + "\n\n…（内容已截断）"
    return text


# ── 各格式提取实现 ──────────────────────────────────────────


def _extract_pdf(data: bytes) -> str:
    import pymupdf  # lazy import

    text_parts: list[str] = []
    with pymupdf.open(stream=data, filetype="pdf") as doc:
        for page in doc:
            page_text = page.get_text()
            if page_text.strip():
                text_parts.append(page_text)
    return "\n".join(text_parts)


def _extract_docx(data: bytes) -> str:
    from docx import Document  # python-docx

    doc = Document(io.BytesIO(data))
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    return "\n".join(paragraphs)


def _extract_pptx(data: bytes) -> str:
    from pptx import Presentation  # python-pptx

    prs = Presentation(io.BytesIO(data))
    text_parts: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    line = paragraph.text.strip()
                    if line:
                        text_parts.append(line)
    return "\n".join(text_parts)


def _extract_txt(data: bytes) -> str:
    # 尝试 UTF-8，失败则用 GBK（适配中文 Windows 文件）
    for encoding in ("utf-8", "gbk", "latin-1"):
        try:
            return data.decode(encoding)
        except (UnicodeDecodeError, LookupError):
            continue
    return data.decode("utf-8", errors="replace")
